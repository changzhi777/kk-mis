"""Office 核心引擎（CLI + API 共用 service 层）。

6 个核心函数，admin 本地处理（无网络依赖）：
- docx_to_pdf：docx → html（mammoth）→ PDF（weasyprint）
- html_to_pdf：HTML 字符串 → PDF（weasyprint）
- data_to_excel：JSON/列表 → xlsx（openpyxl）
- template_to_pptx：模板 + 数据 → PPT（python-pptx）
- fill_form：docx 模板 + 变量 → 填充文档（docxtpl）
- batch_process：文件夹级批量（asyncio + glob）

**2026-07-17 OFFICE-ENGINE-SANDBOX 修复**：
所有函数现在接受 ``OfficeWorkspace`` + 相对路径，输出统一写到
``workspace.temp_path()`` 临时文件再返回 Path。路径遍历/绝对路径
由 workspace.resolve() 拦截；临时文件 mtime > TTL 由 lifespan 后台
任务清理。

设计原则（SOLID/KISS/DRY）：
- 每个函数单一职责，输入/输出清晰
- CLI 和 API 调同一函数（零重复）
- 异常统一 OfficeEngineError
"""
from __future__ import annotations

import asyncio
import io
import json
import logging
import shutil
from pathlib import Path
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from .workspace import OfficeWorkspace

logger = logging.getLogger(__name__)


class OfficeEngineError(Exception):
    """office 引擎处理失败。"""


# ── PDF ──────────────────────────────────────────────────────

def html_to_pdf(
    workspace: "OfficeWorkspace",
    html: str,
    output_name: str | None = None,
) -> Path:
    """HTML 字符串 → PDF（weasyprint），输出在沙箱内临时文件。

    Args:
        workspace: 统一 workspace 沙箱
        html: HTML 内容
        output_name: 可选输出文件名（含后缀，如 ``report.pdf``），

                    **仅作为日志/审计标识；实际写到哪里由 ``output`` 决定**。
                    若提供，必须是相对路径（workspace.resolve 校验）。
    Returns:
        沙箱内 PDF 临时文件路径
    """
    try:
        from weasyprint import HTML
    except ImportError as e:
        raise OfficeEngineError("weasyprint 未安装") from e
    output = workspace.temp_path(suffix=".pdf")
    try:
        HTML(string=html).write_pdf(str(output))
    except OSError as e:
        # weasyprint 依赖 GTK 系统库（libgobject-2.0-0 等）；缺库时抛 OSError
        # 包成 OfficeEngineError 让上层统一处理（→ HTTP 400 而非 500）
        raise OfficeEngineError(f"weasyprint 系统依赖缺失: {e}") from e
    if output_name:
        # 仅做路径校验 + 警告，不覆盖实际 output（避免恶意 output_name 越界）
        try:
            workspace.resolve(output_name)
        except ValueError as e:
            logger.warning(f"html_to_pdf: rejected output_name {output_name!r}: {e}")
    return output


def docx_to_pdf(
    workspace: "OfficeWorkspace",
    input_name: str,
    output_name: str | None = None,
) -> Path:
    """docx → html（mammoth）→ PDF（weasyprint），输入在沙箱内，输出临时文件。

    两步转换：mammoth 提取 html → weasyprint 渲染 PDF。
    格式还原度不如 libreoffice，但纯 Python 无系统依赖。
    """
    # 路径遏制先于依赖检查（缺 mammoth 不该泄露路径是否合法）
    src = workspace.resolve(input_name, must_exist=True)
    try:
        import mammoth
    except ImportError as e:
        raise OfficeEngineError("mammoth 未安装") from e
    with open(src, "rb") as f:
        result = mammoth.convert_to_html(f)
    html = f"<html><body>{result.value}</body></html>"
    return html_to_pdf(workspace, html, output_name=output_name)


# ── Excel ────────────────────────────────────────────────────

def data_to_excel(
    workspace: "OfficeWorkspace",
    data: list[list[Any]] | list[dict[str, Any]],
    output_name: str | None = None,
    headers: list[str] | None = None,
    sheet_name: str = "Sheet1",
) -> Path:
    """JSON/列表 → xlsx（openpyxl），输出在沙箱内临时文件。

    支持两种输入：
    - list[list]：二维数组，headers 可选
    - list[dict]：字典列表，keys 自动提取为 headers
    """
    try:
        from openpyxl import Workbook
    except ImportError as e:
        raise OfficeEngineError("openpyxl 未安装") from e
    output = workspace.temp_path(suffix=".xlsx")
    if output_name:
        try:
            workspace.resolve(output_name)
        except ValueError as e:
            logger.warning(f"data_to_excel: rejected output_name {output_name!r}: {e}")
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    # 统一为 list[list] 格式
    rows: list[list[Any]] = []
    if data and isinstance(data[0], dict):
        keys = headers or list(data[0].keys())
        rows.append(keys)
        for item in data:
            rows.append([item.get(k, "") for k in keys])
    else:
        if headers:
            rows.append(headers)
        rows.extend(data)
    for row in rows:
        ws.append(row)
    wb.save(str(output))
    return output


# ── PPT ──────────────────────────────────────────────────────

def template_to_pptx(
    workspace: "OfficeWorkspace",
    slides: list[dict[str, Any]],
    output_name: str | None = None,
    template: str | None = None,
) -> Path:
    """数据 → PPT（python-pptx），输出在沙箱内临时文件。

    每条 slide dict 含：title / content / layout（可选，默认 1=标题+内容）
    template 可选：加载已有 pptx 作为基础模板（沙箱内相对路径）。
    """
    # 路径遏制先于依赖检查；template=None 时跳过
    template_path = (
        workspace.resolve(template, must_exist=True) if template else None
    )
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise OfficeEngineError("python-pptx 未安装") from e
    output = workspace.temp_path(suffix=".pptx")
    if output_name:
        try:
            workspace.resolve(output_name)
        except ValueError as e:
            logger.warning(f"template_to_pptx: rejected output_name {output_name!r}: {e}")
    prs = Presentation(str(template_path)) if template_path else Presentation()
    for slide_data in slides:
        layout_idx = slide_data.get("layout", 1)
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        if "title" in slide_data:
            slide.shapes.title.text = str(slide_data["title"])
        if "content" in slide_data:
            body = slide.placeholders[1]
            body.text = str(slide_data["content"])
    prs.save(str(output))
    return output


# ── 表单填写 ──────────────────────────────────────────────────

def fill_form(
    workspace: "OfficeWorkspace",
    template: str,
    data: dict[str, Any],
    output_name: str | None = None,
) -> Path:
    """docx 模板 + 变量 → 填充文档（docxtpl），输出在沙箱内临时文件。

    模板内用 {{var}} 占位符，data 提供 key-value。
    """
    # 路径遏制先于依赖检查（缺 docxtpl 不该泄露路径是否合法）
    template_path = workspace.resolve(template, must_exist=True)
    try:
        from docxtpl import DocxTemplate
    except ImportError as e:
        raise OfficeEngineError("docxtpl 未安装") from e
    output = workspace.temp_path(suffix=".docx")
    if output_name:
        try:
            workspace.resolve(output_name)
        except ValueError as e:
            logger.warning(f"fill_form: rejected output_name {output_name!r}: {e}")
    doc = DocxTemplate(str(template_path))
    doc.render(data)
    doc.save(str(output))
    return output


# ── 批量处理 ──────────────────────────────────────────────────

async def batch_process(
    workspace: "OfficeWorkspace",
    input_dir: str,
    operation: str,
    output_dir: str | None = None,
    pattern: str = "*",
) -> list[Path]:
    """文件夹级批量处理，全部在 workspace 沙箱内。

    operation 支持：pdf（docx→PDF）/ form（需 data.json 同名文件）/ copy
    遍历 ``workspace.resolve(input_dir)`` 下匹配 ``pattern`` 的文件，
    逐个调用对应函数，输出写到沙箱内的 ``output_dir``（沙箱校验）。

    Args:
        workspace: 统一 workspace 沙箱
        input_dir: 输入目录相对路径（沙箱校验）
        operation: pdf | form | copy
        output_dir: 输出目录相对路径；缺省在 ``input_dir/output_<op>`` 下
        pattern: glob 匹配模式

    Returns:
        输出文件 Path 列表（全部在沙箱内）
    """
    try:
        in_dir = workspace.resolve(input_dir, must_exist=True)
        if output_dir:
            out_dir = workspace.resolve(output_dir)
        else:
            # 同沙箱内子目录：input_dir/output_<op>
            out_dir = workspace.resolve(f"{input_dir.rstrip('/')}/output_{operation}")
    except ValueError as e:
        # workspace.resolve 抛 ValueError（路径逃出沙箱）→ 转 OfficeEngineError
        # 让路由 try/except 统一捕获并返回 HTTP 400（而非 500）
        raise OfficeEngineError(f"路径不在 workspace 沙箱内: {e}") from e
    out_dir.mkdir(parents=True, exist_ok=True)
    files = sorted(in_dir.glob(pattern))
    results: list[Path] = []
    for f in files:
        if not f.is_file():
            continue
        try:
            out_base = out_dir / f"{f.stem}_converted"
            if operation == "pdf":
                if f.suffix.lower() == ".docx":
                    # 单文件复用 docx_to_pdf 沙箱逻辑
                    src_rel = str(f.relative_to(workspace.root))
                    result = docx_to_pdf(workspace, src_rel)
                    # 移入 out_dir（沙箱内 mv 仍是沙箱内）
                    target = out_base.with_suffix(".pdf")
                    shutil.move(str(result), str(target))
                    results.append(target)
            elif operation == "form":
                data_file = f.with_suffix(".json")
                if data_file.exists():
                    data = json.loads(data_file.read_text(encoding="utf-8"))
                    tpl_rel = str(f.relative_to(workspace.root))
                    result = fill_form(workspace, tpl_rel, data)
                    target = out_base.with_suffix(".docx")
                    shutil.move(str(result), str(target))
                    results.append(target)
            elif operation == "copy":
                dest = out_dir / f.name
                shutil.copy2(f, dest)
                results.append(dest)
            else:
                raise OfficeEngineError(f"未知操作: {operation}")
        except Exception as e:
            logger.warning("批量处理 %s 失败: %s", f, e)
    return results